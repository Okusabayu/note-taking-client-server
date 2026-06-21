"""
Script untuk generate file PowerPoint (.pptx) presentasi akademik
Tugas Akhir Client Server Programming
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================================================
# KONFIGURASI WARNA
# ==================================================
COLOR_PRIMARY = RGBColor(79, 70, 229)      # Indigo
COLOR_PRIMARY_DARK = RGBColor(67, 56, 202)
COLOR_PURPLE = RGBColor(124, 58, 237)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(30, 41, 59)
COLOR_GRAY = RGBColor(100, 116, 139)
COLOR_LIGHT_BG = RGBColor(241, 245, 249)
COLOR_LIGHT_GRAY = RGBColor(226, 232, 240)
COLOR_SUCCESS = RGBColor(34, 197, 94)
COLOR_DANGER = RGBColor(239, 68, 68)
COLOR_WARNING = RGBColor(245, 158, 11)
COLOR_DARK_BG = RGBColor(15, 23, 42)
COLOR_ACCENT = RGBColor(99, 102, 241)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ==================================================
# HELPER FUNCTIONS
# ==================================================
def add_gradient_bg(slide, color1=COLOR_DARK_BG, color2=COLOR_PRIMARY):
    """Add a solid dark background rectangle."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color1
    bg.line.fill.background()


def add_light_bg(slide):
    """Add a light background."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_WHITE
    bg.line.fill.background()


def add_header_bar(slide, title_text, subtitle_text=None):
    """Add a colored header bar at top."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Calibri"

    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(200, 200, 255)
        p2.font.name = "Calibri"


def add_text_block(slide, left, top, width, height, text, font_size=14, color=COLOR_BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text block."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=COLOR_BLACK):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return tf


def add_table(slide, left, top, width, height, data, col_widths=None):
    """Add a table. data[0] is the header row."""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLOR_WHITE
                else:
                    paragraph.font.color.rgb = COLOR_BLACK

            # Header row coloring
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_PRIMARY
            else:
                cell.fill.solid()
                if row_idx % 2 == 0:
                    cell.fill.fore_color.rgb = COLOR_LIGHT_BG
                else:
                    cell.fill.fore_color.rgb = COLOR_WHITE

    return table


def add_accent_line(slide, left, top, width):
    """Add a small accent line."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()


def add_code_block(slide, left, top, width, height, code_text):
    """Add a code block with dark background."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(30, 30, 46)
    bg.line.fill.background()

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(166, 227, 161)
    p.font.name = "Consolas"
    return tf


def add_card(slide, left, top, width, height, title, content, icon=""):
    """Add a card-style box."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_BG
    card.line.color.rgb = COLOR_LIGHT_GRAY
    card.line.width = Pt(1)

    # Title
    add_text_block(slide, left + Inches(0.2), top + Inches(0.1),
                   width - Inches(0.4), Inches(0.4),
                   f"{icon} {title}" if icon else title,
                   font_size=13, bold=True, color=COLOR_PRIMARY)

    # Content
    add_text_block(slide, left + Inches(0.2), top + Inches(0.45),
                   width - Inches(0.4), height - Inches(0.55),
                   content, font_size=11, color=COLOR_GRAY)


# ==================================================
# SLIDE 1: HALAMAN JUDUL
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_gradient_bg(slide)

# Accent line top
accent = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.06)
)
accent.fill.solid()
accent.fill.fore_color.rgb = COLOR_ACCENT
accent.line.fill.background()

# Emoji / icon
add_text_block(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(0.8),
               "📝", font_size=48, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)

# Main title
add_text_block(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2),
               "Analisis dan Implementasi Sederhana\nAplikasi Note Taking Berbasis Client-Server",
               font_size=32, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle line
accent2 = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.4), Inches(3.3), Inches(0.04)
)
accent2.fill.solid()
accent2.fill.fore_color.rgb = COLOR_ACCENT
accent2.line.fill.background()

add_text_block(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.4),
               "Tugas Akhir Mata Kuliah Client Server Programming",
               font_size=16, color=RGBColor(165, 180, 252), alignment=PP_ALIGN.CENTER)

# Info
info_text = "Disusun Oleh: [Nama Mahasiswa]  •  NIM: [NIM]"
add_text_block(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.4),
               info_text, font_size=14, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

add_text_block(slide, Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.4),
               "Dosen Pengampu: [Nama Dosen]",
               font_size=14, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)

add_text_block(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.4),
               "[Program Studi]  •  [Universitas]  •  Semester Genap 2025/2026",
               font_size=12, color=RGBColor(80, 90, 110), alignment=PP_ALIGN.CENTER)


# ==================================================
# SLIDE 2: PENDAHULUAN
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Pendahuluan", "Latar Belakang, Rumusan Masalah, dan Tujuan")

# Latar Belakang
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
               "Latar Belakang", font_size=18, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(0.8), Inches(2.0), Inches(1.5))

add_text_block(slide, Inches(0.8), Inches(2.15), Inches(5.5), Inches(1.5),
               "Perkembangan teknologi informasi telah mengubah cara manusia mengelola informasi. "
               "Pencatatan digital (note-taking) menjadi kebutuhan fundamental yang memungkinkan "
               "pengguna menyimpan dan mengakses catatan dari berbagai perangkat. "
               "Arsitektur client-server menjadi fondasi utama aplikasi modern seperti Google Keep.",
               font_size=13, color=COLOR_BLACK)

# Rumusan Masalah
add_text_block(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.4),
               "Rumusan Masalah", font_size=18, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(0.8), Inches(4.2), Inches(1.5))

add_bullet_list(slide, Inches(0.8), Inches(4.35), Inches(5.5), Inches(2.0), [
    "Bagaimana arsitektur client-server diterapkan pada Google Keep?",
    "Bagaimana mengimplementasikan note-taking app dengan Flask + SQLite?",
    "Bagaimana perbandingan implementasi sederhana vs sistem produksi?"
], font_size=12)

# Tujuan (right side)
add_text_block(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
               "Tujuan Penelitian", font_size=18, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(7), Inches(2.0), Inches(1.5))

add_bullet_list(slide, Inches(7), Inches(2.15), Inches(5.5), Inches(4.0), [
    "Menganalisis arsitektur client-server pada Google Keep",
    "Mengimplementasikan aplikasi note-taking dengan Flask & SQLite",
    "Mengevaluasi aspek komunikasi, konkurensi, dan keamanan"
], font_size=13)

# Tech stack cards
add_card(slide, Inches(7), Inches(4.0), Inches(2.5), Inches(1.2),
         "Backend", "Python + Flask\nSQLite Database", "🐍")
add_card(slide, Inches(9.7), Inches(4.0), Inches(2.5), Inches(1.2),
         "Frontend", "HTML + CSS\nJavaScript (Fetch API)", "🌐")


# ==================================================
# SLIDE 3: TINJAUAN PUSTAKA
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Tinjauan Pustaka", "Dasar Teori dan Teknologi")

# REST principles table
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
               "REST (Representational State Transfer)", font_size=16, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(0.8), Inches(2.0), Inches(2))

add_text_block(slide, Inches(0.8), Inches(2.15), Inches(5.5), Inches(0.6),
               "Gaya arsitektur untuk API web yang dikemukakan oleh Roy Fielding (2000). "
               "Menggunakan HTTP methods standar untuk operasi CRUD.",
               font_size=12, color=COLOR_GRAY)

rest_data = [
    ["Prinsip", "Penjelasan"],
    ["Stateless", "Setiap request berdiri sendiri"],
    ["Client-Server", "Pemisahan antarmuka dan data"],
    ["Uniform Interface", "HTTP methods standar"],
    ["Resource-Based", "Entitas diakses via URI unik"],
]
add_table(slide, Inches(0.8), Inches(2.9), Inches(5.5), Inches(2.0), rest_data)

# Tech stack
add_text_block(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
               "Teknologi yang Digunakan", font_size=16, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(7), Inches(2.0), Inches(2))

tech_data = [
    ["Teknologi", "Peran", "Versi"],
    ["Python", "Bahasa pemrograman", "3.x"],
    ["Flask", "Web micro-framework", "3.1.3"],
    ["SQLite", "Database relasional", "Built-in"],
    ["JavaScript", "Client-side logic", "ES6+"],
    ["HTML/CSS", "Antarmuka pengguna", "HTML5"],
]
add_table(slide, Inches(7), Inches(2.15), Inches(5.5), Inches(2.8), tech_data)

# Client-Server definition
add_text_block(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.4),
               "Model Client-Server", font_size=16, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(0.8), Inches(5.7), Inches(1.5))
add_text_block(slide, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.8),
               "Model arsitektur di mana client mengirim permintaan (request) dan server memproses serta mengembalikan respons (response). "
               "Merupakan fondasi dari sebagian besar aplikasi web modern (Tanenbaum & Van Steen, 2017).",
               font_size=12, color=COLOR_GRAY)


# ==================================================
# SLIDE 4: METODOLOGI PENELITIAN
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Metodologi Penelitian", "Research and Development (R&D)")

# Flow steps as cards
steps = [
    ("1", "Studi Literatur", "Mempelajari arsitektur client-server, REST API, Flask", "📚"),
    ("2", "Analisis Sistem", "Menganalisis Google Keep sebagai studi kasus", "🔍"),
    ("3", "Perancangan", "Merancang arsitektur, database, dan endpoint API", "📐"),
    ("4", "Implementasi", "Coding backend (Flask + SQLite) dan frontend (HTML/JS)", "💻"),
    ("5", "Pengujian", "Testing endpoint API menggunakan Postman", "🧪"),
    ("6", "Evaluasi", "Membandingkan implementasi vs sistem produksi", "📊"),
]

for i, (num, title, desc, icon) in enumerate(steps):
    row = i // 3
    col = i % 3
    left = Inches(0.8) + col * Inches(4.1)
    top = Inches(1.8) + row * Inches(2.4)

    # Card background
    card_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.0)
    )
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = COLOR_LIGHT_BG
    card_shape.line.color.rgb = COLOR_LIGHT_GRAY
    card_shape.line.width = Pt(1)

    # Number badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15), Inches(0.45), Inches(0.45)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_PRIMARY
    badge.line.fill.background()
    add_text_block(slide, left + Inches(0.15), top + Inches(0.18), Inches(0.45), Inches(0.4),
                   num, font_size=16, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Title & icon
    add_text_block(slide, left + Inches(0.7), top + Inches(0.2), Inches(2.8), Inches(0.35),
                   f"{icon} {title}", font_size=15, bold=True, color=COLOR_BLACK)

    # Description
    add_text_block(slide, left + Inches(0.2), top + Inches(0.7), Inches(3.3), Inches(1.1),
                   desc, font_size=12, color=COLOR_GRAY)

    # Arrow between cards (except last in row)
    if col < 2 and i < 5:
        add_text_block(slide, left + Inches(3.7), top + Inches(0.7), Inches(0.4), Inches(0.4),
                       "→", font_size=20, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)


# ==================================================
# SLIDE 5: ANALISIS GOOGLE KEEP
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Analisis Google Keep", "Studi Kasus Aplikasi Pencatatan Digital Skala Produksi")

add_text_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.7),
               "Google Keep adalah layanan pencatatan digital oleh Google, tersedia sebagai "
               "web app, Android, dan iOS. Diluncurkan 2013 dengan infrastruktur Google Cloud.",
               font_size=12, color=COLOR_GRAY)

# Features table
add_text_block(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(0.3),
               "Fitur Utama", font_size=15, bold=True, color=COLOR_PRIMARY)

feat_data = [
    ["Kategori", "Fitur"],
    ["CRUD", "Buat, baca, edit, hapus catatan"],
    ["Organisasi", "Label, warna, pin, arsip"],
    ["Kolaborasi", "Berbagi catatan, real-time editing"],
    ["Media", "Gambar, audio, drawing"],
    ["Sinkronisasi", "Cross-device, offline-first"],
    ["AI", "OCR, voice-to-text"],
]
add_table(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(3.5), feat_data)

# Architecture (simplified text version)
add_text_block(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.3),
               "Arsitektur Google Keep (Estimasi)", font_size=15, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(7), Inches(1.95), Inches(2))

arch_layers = [
    ("🌐 Client Layer", "Web App, Android, iOS"),
    ("🔀 API Gateway", "Load Balancer + Auth (OAuth 2.0)"),
    ("⚙️ Service Layer", "Notes, Media, Sync, Search (Microservices)"),
    ("🗄️ Data Layer", "Cloud Spanner, Redis Cache, Cloud Storage"),
]
for i, (layer, desc) in enumerate(arch_layers):
    top = Inches(2.15) + i * Inches(1.15)
    add_card(slide, Inches(7), top, Inches(5.5), Inches(0.95), layer, desc)

    # Arrow
    if i < len(arch_layers) - 1:
        add_text_block(slide, Inches(9.5), top + Inches(0.9), Inches(0.5), Inches(0.3),
                       "↓", font_size=14, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)


# ==================================================
# SLIDE 6: ARSITEKTUR CLIENT-SERVER
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Arsitektur Client-Server", "Perbandingan Implementasi Sederhana vs Google Keep")

# Left: Our implementation
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.4),
               "🛠️ Implementasi Sederhana (Monolitik 2-Tier)", font_size=15, bold=True, color=COLOR_PRIMARY)

impl_items = [
    ("🌐", "Browser", "HTML + CSS + JavaScript"),
    ("↕️", "", "HTTP Request / JSON Response"),
    ("🐍", "Flask Server", "REST API + Template Engine"),
    ("↕️", "", "SQL Query / Result Set"),
    ("🗄️", "SQLite", "File-based database"),
]
for i, (icon, title, desc) in enumerate(impl_items):
    top = Inches(2.1) + i * Inches(0.7)
    if title:
        add_card(slide, Inches(1.2), top, Inches(4.5), Inches(0.6), f"{icon} {title}", desc)
    else:
        add_text_block(slide, Inches(3), top, Inches(1), Inches(0.5),
                       icon, font_size=16, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)

# Right: Google Keep
add_text_block(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.4),
               "☁️ Google Keep (Microservices N-Tier)", font_size=15, bold=True, color=COLOR_PRIMARY)

gk_items = [
    ("📱", "Multi-Platform Client", "Web, Android, iOS"),
    ("↕️", "", "HTTPS + OAuth 2.0"),
    ("🔀", "API Gateway", "Load Balancer + Auth"),
    ("↕️", "", "gRPC / Protobuf"),
    ("⚙️", "Microservices", "Notes, Media, Sync, Search"),
    ("↕️", "", "Distributed Query"),
    ("🗄️", "Cloud Database", "Spanner + Redis + Storage"),
]
for i, (icon, title, desc) in enumerate(gk_items):
    top = Inches(2.1) + i * Inches(0.55)
    if title:
        add_card(slide, Inches(7.4), top, Inches(4.5), Inches(0.5), f"{icon} {title}", desc)
    else:
        add_text_block(slide, Inches(9.2), top, Inches(1), Inches(0.4),
                       icon, font_size=14, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)

# Comparison table at bottom
comp_data = [
    ["Aspek", "Implementasi Sederhana", "Google Keep"],
    ["Arsitektur", "Monolitik (2-tier)", "Microservices (N-tier)"],
    ["Database", "SQLite (file-based)", "Cloud Spanner / Bigtable"],
    ["Skalabilitas", "Single user", "Jutaan concurrent users"],
    ["Deployment", "Localhost", "Google Cloud Platform"],
]
add_table(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.5), comp_data)


# ==================================================
# SLIDE 7: PROTOKOL KOMUNIKASI DAN FORMAT DATA
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Protokol Komunikasi dan Format Data", "HTTP Methods, JSON, dan Endpoint REST")

# Endpoint table
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.3),
               "Endpoint REST API", font_size=16, bold=True, color=COLOR_PRIMARY)

ep_data = [
    ["Method", "Endpoint", "Fungsi", "Request Body", "Response"],
    ["GET", "/notes", "Ambil semua catatan", "—", "JSON array (200)"],
    ["POST", "/notes", "Buat catatan baru", "{title, content}", "message (201)"],
    ["PUT", "/notes/<id>", "Update catatan", "{title, content}", "message (200)"],
    ["DELETE", "/notes/<id>", "Hapus catatan", "—", "message (200)"],
]
add_table(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.2), ep_data)

# Example request/response
add_text_block(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.3),
               "Contoh Request (POST)", font_size=14, bold=True, color=COLOR_PRIMARY)

add_code_block(slide, Inches(0.8), Inches(4.9), Inches(5.5), Inches(1.6),
               'POST /notes HTTP/1.1\n'
               'Content-Type: application/json\n\n'
               '{\n'
               '    "title": "Catatan Pertama",\n'
               '    "content": "Isi catatan"\n'
               '}')

# Comparison
add_text_block(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(0.3),
               "Perbandingan Protokol", font_size=14, bold=True, color=COLOR_PRIMARY)

proto_data = [
    ["Aspek", "Implementasi", "Google Keep"],
    ["Protokol", "HTTP", "HTTPS + TLS 1.3"],
    ["Format Data", "JSON", "Protobuf + JSON"],
    ["Auth", "Tidak ada", "OAuth 2.0"],
    ["Enkripsi", "Tidak ada", "End-to-end"],
]
add_table(slide, Inches(7), Inches(4.9), Inches(5.5), Inches(1.6), proto_data)


# ==================================================
# SLIDE 8: PENANGANAN KONKURENSI
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Penanganan Konkurensi", "Threading dan Database Locking")

# Left side: explanation
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
               "Mekanisme Konkurensi", font_size=16, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(0.8), Inches(1.95), Inches(1.5))

add_text_block(slide, Inches(0.8), Inches(2.15), Inches(5.5), Inches(1.0),
               "Flask menangani setiap HTTP request dalam thread terpisah. "
               "SQLite menggunakan file-level locking — hanya satu operasi write "
               "yang dapat berjalan pada satu waktu. Cocok untuk traffic rendah.",
               font_size=12, color=COLOR_GRAY)

# Flow visualization
flow_steps = [
    "User 1 → POST /notes → Thread 1",
    "User 2 → POST /notes → Thread 2",
    "Thread 1 → SQLite LOCK → INSERT → OK",
    "Thread 2 → WAIT → SQLite LOCK → INSERT → OK",
]
add_bullet_list(slide, Inches(0.8), Inches(3.3), Inches(5.5), Inches(2.0),
                flow_steps, font_size=12, color=COLOR_BLACK)

# Right side: comparison table
add_text_block(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.3),
               "Perbandingan Mekanisme", font_size=16, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(7), Inches(1.95), Inches(1.5))

conc_data = [
    ["Aspek", "Implementasi", "Google Keep"],
    ["Threading", "1 thread/request", "Thread pool + async"],
    ["DB Locking", "File-level lock", "Row-level lock"],
    ["Connection", "Per-request", "Connection pooling"],
    ["Conflict", "Last-write-wins", "OT / CRDT"],
    ["Caching", "Tidak ada", "Multi-layer cache"],
    ["Max Users", "~10–50", "Jutaan"],
]
add_table(slide, Inches(7), Inches(2.15), Inches(5.5), Inches(3.5), conc_data)

# Bottom note
add_text_block(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6),
               "⚠️ Keterbatasan: SQLite cocok untuk development & single-server. "
               "Untuk produksi, disarankan migrasi ke PostgreSQL atau MySQL yang mendukung row-level locking.",
               font_size=12, bold=True, color=COLOR_WARNING)


# ==================================================
# SLIDE 9: ASPEK KEAMANAN
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Aspek Keamanan", "Perbandingan Lapisan Keamanan")

# Security comparison table
sec_data = [
    ["Lapisan", "Implementasi", "Google Keep"],
    ["Transport", "HTTP (plain text)", "HTTPS + TLS 1.3"],
    ["Autentikasi", "Tidak ada", "OAuth 2.0 + 2FA"],
    ["Otorisasi", "Akses terbuka", "Role-based, per-user"],
    ["SQL Injection", "✅ Parameterized query", "Multi-layer validation"],
    ["XSS", "✅ HTML escaping", "CSP + sanitization"],
    ["CSRF", "Tidak ada", "Token-based"],
    ["Enkripsi Data", "Tidak ada", "At-rest + in-transit"],
]
add_table(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(3.2), sec_data)

# Code examples
add_text_block(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(0.3),
               "✅ Parameterized Query (SQL Injection Safe)", font_size=13, bold=True, color=COLOR_SUCCESS)

add_code_block(slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.2),
               '# Aman - parameterized query\n'
               'conn.execute(\n'
               '    "INSERT INTO notes VALUES (?, ?)",\n'
               '    (title, content)\n'
               ')')

add_text_block(slide, Inches(7), Inches(5.1), Inches(5.5), Inches(0.3),
               "✅ Input Validation + XSS Escaping", font_size=13, bold=True, color=COLOR_SUCCESS)

add_code_block(slide, Inches(7), Inches(5.5), Inches(5.5), Inches(1.2),
               '# Server-side validation\n'
               'if not title or not content:\n'
               '    return jsonify({"error": "..."}), 400\n\n'
               '# Client-side XSS escaping\n'
               'div.textContent = text;  // safe')


# ==================================================
# SLIDE 10: DESAIN IMPLEMENTASI
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Desain Implementasi", "Database, Struktur Folder, dan Alur Sistem")

# ERD
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(3.5), Inches(0.3),
               "Entity Relationship", font_size=15, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(0.8), Inches(1.95), Inches(1.5))

erd_data = [
    ["Kolom", "Tipe", "Keterangan"],
    ["id", "INTEGER", "PRIMARY KEY, AUTO INCREMENT"],
    ["title", "TEXT", "NOT NULL"],
    ["content", "TEXT", "NOT NULL"],
]
add_table(slide, Inches(0.8), Inches(2.15), Inches(3.8), Inches(1.6), erd_data)

# Folder structure
add_text_block(slide, Inches(0.8), Inches(4.2), Inches(3.8), Inches(0.3),
               "Struktur Folder", font_size=15, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(0.8), Inches(4.55), Inches(1.5))

add_code_block(slide, Inches(0.8), Inches(4.7), Inches(3.8), Inches(2.2),
               'note-taking-csp/\n'
               '├── static/\n'
               '│   ├── style.css\n'
               '│   └── script.js\n'
               '├── templates/\n'
               '│   └── index.html\n'
               '├── app.py\n'
               '├── database.db\n'
               '├── init_db.py\n'
               '└── requirements.txt')

# System flow (right side)
add_text_block(slide, Inches(5.4), Inches(1.6), Inches(7), Inches(0.3),
               "Alur Sistem (Flowchart)", font_size=15, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(5.4), Inches(1.95), Inches(1.5))

flow_items = [
    ("1️⃣", "User membuka browser → GET /"),
    ("2️⃣", "Flask mengirim index.html + CSS + JS"),
    ("3️⃣", "JavaScript: fetch GET /notes"),
    ("4️⃣", "Flask query SQLite → return JSON"),
    ("5️⃣", "JS render daftar catatan ke DOM"),
    ("6️⃣", "User klik Tambah → POST /notes"),
    ("7️⃣", "Flask INSERT ke SQLite → 201"),
    ("8️⃣", "JS refresh daftar otomatis"),
    ("9️⃣", "Edit → PUT /notes/id | Hapus → DELETE /notes/id"),
]

for i, (num, text) in enumerate(flow_items):
    top = Inches(2.15) + i * Inches(0.52)
    add_text_block(slide, Inches(5.4), top, Inches(7.2), Inches(0.45),
                   f"{num}  {text}", font_size=12, color=COLOR_BLACK)


# ==================================================
# SLIDE 11: IMPLEMENTASI PROGRAM
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Implementasi Program", "Kode Backend (Flask) dan Frontend (JavaScript)")

# Backend
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
               "🐍 Backend — Flask REST API", font_size=15, bold=True, color=COLOR_PRIMARY)

add_code_block(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.5),
               'def get_db_connection():\n'
               '    conn = sqlite3.connect("database.db")\n'
               '    conn.row_factory = sqlite3.Row\n'
               '    return conn\n\n'
               '@app.route("/")\n'
               'def home():\n'
               '    return render_template("index.html")')

add_code_block(slide, Inches(0.8), Inches(3.7), Inches(5.5), Inches(1.8),
               '@app.route("/notes", methods=["GET"])\n'
               'def get_notes():\n'
               '    conn = get_db_connection()\n'
               '    notes = conn.execute(\n'
               '        "SELECT * FROM notes"\n'
               '    ).fetchall()\n'
               '    conn.close()\n'
               '    return jsonify([...])  # JSON response')

# Frontend
add_text_block(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.3),
               "🌐 Frontend — JavaScript Fetch API", font_size=15, bold=True, color=COLOR_PRIMARY)

add_code_block(slide, Inches(7), Inches(2.0), Inches(5.5), Inches(1.5),
               '// Mengambil catatan (GET)\n'
               'async function loadNotes() {\n'
               '    const res = await fetch("/notes");\n'
               '    const notes = await res.json();\n'
               '    // Render ke DOM...\n'
               '}')

add_code_block(slide, Inches(7), Inches(3.7), Inches(5.5), Inches(1.8),
               '// Menambah catatan (POST)\n'
               'async function createNote(title, content) {\n'
               '    await fetch("/notes", {\n'
               '        method: "POST",\n'
               '        headers: {\n'
               '          "Content-Type": "application/json"\n'
               '        },\n'
               '        body: JSON.stringify({title, content})\n'
               '    });\n'
               '}')

# Bottom note
add_text_block(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.8),
               "📌 Poin Penting: Backend menyediakan REST API dan melayani halaman HTML via Jinja2 template engine. "
               "Frontend menggunakan vanilla JavaScript (tanpa framework) dengan Fetch API untuk komunikasi asynchronous. "
               "Semua operasi CRUD dilakukan tanpa reload halaman (Single Page Application behavior).",
               font_size=12, bold=False, color=COLOR_GRAY)


# ==================================================
# SLIDE 12: HASIL IMPLEMENTASI
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Hasil Implementasi", "Tampilan Aplikasi dan Alur Komunikasi")

# UI Components
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
               "Komponen Antarmuka", font_size=15, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(0.8), Inches(1.95), Inches(1.5))

ui_data = [
    ["Komponen", "Deskripsi"],
    ["Header", "Judul app + gradient indigo-ungu"],
    ["Form Input", "Field judul, isi catatan, tombol Tambah"],
    ["Daftar Catatan", "Card list + tombol Edit & Hapus"],
    ["Toast", "Notifikasi sukses/error (pojok kanan bawah)"],
    ["Responsive", "Menyesuaikan desktop & mobile"],
]
add_table(slide, Inches(0.8), Inches(2.15), Inches(5.5), Inches(2.5), ui_data)

# Communication flow
add_text_block(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.3),
               "Alur Komunikasi", font_size=15, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(7), Inches(1.95), Inches(1.5))

comm_flow = [
    "🌐 Browser → GET / → Flask → index.html",
    "📜 script.js → GET /notes → Flask → SQLite → JSON",
    "➕ Tambah → POST /notes → Flask → INSERT → 201",
    "✏️ Edit → PUT /notes/id → Flask → UPDATE → 200",
    "🗑️ Hapus → DELETE /notes/id → Flask → DELETE → 200",
    "🔄 Setiap aksi → auto refresh daftar catatan",
]
add_bullet_list(slide, Inches(7), Inches(2.15), Inches(5.5), Inches(3.0),
                comm_flow, font_size=12)

# Bottom: screenshot placeholder
screenshot_bg = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8)
)
screenshot_bg.fill.solid()
screenshot_bg.fill.fore_color.rgb = COLOR_LIGHT_BG
screenshot_bg.line.color.rgb = COLOR_LIGHT_GRAY
screenshot_bg.line.width = Pt(1)

add_text_block(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.0),
               "📸 [Masukkan screenshot aplikasi di sini]\n"
               "Ambil screenshot dari http://127.0.0.1:5000 saat aplikasi berjalan",
               font_size=14, color=COLOR_GRAY, alignment=PP_ALIGN.CENTER)


# ==================================================
# SLIDE 13: PENGUJIAN DAN EVALUASI
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Pengujian dan Evaluasi", "Hasil Testing API (Postman) dan Frontend")

# API Testing
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.3),
               "Pengujian API dengan Postman", font_size=15, bold=True, color=COLOR_PRIMARY)

api_test = [
    ["No", "Test Case", "Method", "Endpoint", "Expected", "Result"],
    ["1", "Ambil semua catatan", "GET", "/notes", "200 + JSON array", "✅ Pass"],
    ["2", "Tambah catatan valid", "POST", "/notes", "201 Created", "✅ Pass"],
    ["3", "Tambah tanpa title", "POST", "/notes", "400 Error", "✅ Pass"],
    ["4", "Edit catatan ada", "PUT", "/notes/1", "200 OK", "✅ Pass"],
    ["5", "Edit catatan tidak ada", "PUT", "/notes/999", "404 Not Found", "✅ Pass"],
    ["6", "Hapus catatan ada", "DELETE", "/notes/1", "200 OK", "✅ Pass"],
    ["7", "Hapus catatan tidak ada", "DELETE", "/notes/999", "404 Not Found", "✅ Pass"],
]
add_table(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.8), api_test)

# Frontend Testing
add_text_block(slide, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.3),
               "Pengujian Frontend", font_size=15, bold=True, color=COLOR_PRIMARY)

fe_test = [
    ["No", "Test Case", "Aksi", "Expected", "Result"],
    ["1", "Halaman dimuat", "Buka localhost:5000", "Tampil + daftar catatan", "✅ Pass"],
    ["2", "Tambah catatan", "Isi form → Tambah", "Catatan baru muncul", "✅ Pass"],
    ["3", "Edit catatan", "Edit → ubah → Simpan", "Data berubah", "✅ Pass"],
    ["4", "Hapus catatan", "Hapus → konfirmasi", "Catatan hilang", "✅ Pass"],
    ["5", "Validasi kosong", "Submit kosong", "Tidak terkirim", "✅ Pass"],
]
add_table(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.6), fe_test)


# ==================================================
# SLIDE 14: PEMBAHASAN
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Pembahasan", "Kelebihan, Keterbatasan, dan Rekomendasi")

# Kelebihan
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
               "✅ Kelebihan Implementasi", font_size=15, bold=True, color=COLOR_SUCCESS)
add_accent_line(slide, Inches(0.8), Inches(1.95), Inches(1.5))

pros = [
    "Sederhana dan mudah dipahami — cocok untuk pembelajaran",
    "Full REST compliance — standar HTTP methods & status codes",
    "Separation of concerns — backend API & frontend JS terpisah",
    "Tanpa framework frontend — pemahaman fundamental",
    "SQL Injection safe — parameterized query diterapkan",
]
add_bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.5), Inches(2.5),
                pros, font_size=12)

# Keterbatasan & Rekomendasi
add_text_block(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.3),
               "⚠️ Keterbatasan & Rekomendasi", font_size=15, bold=True, color=COLOR_WARNING)
add_accent_line(slide, Inches(7), Inches(1.95), Inches(1.5))

limit_data = [
    ["Keterbatasan", "Rekomendasi"],
    ["Tidak ada autentikasi", "JWT / Flask-Login"],
    ["SQLite tidak scalable", "Migrasi ke PostgreSQL"],
    ["HTTP tanpa enkripsi", "HTTPS (Let's Encrypt)"],
    ["Tidak ada pagination", "Limit-offset query"],
    ["Single server", "Deploy ke cloud"],
    ["Tidak ada timestamp", "Tambah created_at"],
]
add_table(slide, Inches(7), Inches(2.15), Inches(5.5), Inches(3.0), limit_data)

# Lessons learned
add_text_block(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.3),
               "📌 Lessons Learned", font_size=15, bold=True, color=COLOR_PRIMARY)

add_text_block(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0),
               "Implementasi ini membuktikan bahwa prinsip-prinsip client-server programming dapat diterapkan "
               "secara nyata dengan teknologi open-source. Meskipun sederhana, arsitektur yang dibangun sudah "
               "mencerminkan pola yang digunakan dalam sistem produksi — hanya berbeda pada skala dan kompleksitas.",
               font_size=13, color=COLOR_GRAY)


# ==================================================
# SLIDE 15: KESIMPULAN DAN SARAN
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Kesimpulan dan Saran")

# Kesimpulan
add_text_block(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
               "Kesimpulan", font_size=18, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(0.8), Inches(2.0), Inches(1.5))

conclusions = [
    "Arsitektur client-server berhasil diterapkan menggunakan Flask (server) dan browser (client) dengan komunikasi REST API + JSON.",
    "Google Keep menggunakan arsitektur yang lebih kompleks (microservices, OAuth 2.0), namun fondasi dasarnya tetap sama — request/response pattern.",
    "Empat endpoint CRUD (GET, POST, PUT, DELETE) berhasil diimplementasikan dan diuji via Postman dan web interface.",
    "Aspek keamanan dasar (parameterized query, input validation, XSS escaping) telah diterapkan pada implementasi.",
]
add_bullet_list(slide, Inches(0.8), Inches(2.15), Inches(5.5), Inches(4.0),
                conclusions, font_size=12)

# Saran
add_text_block(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(0.3),
               "Saran Pengembangan", font_size=18, bold=True, color=COLOR_PRIMARY)
add_accent_line(slide, Inches(7), Inches(2.0), Inches(1.5))

suggestions = [
    "Menambahkan sistem autentikasi (JWT / session-based) untuk isolasi data per-user.",
    "Migrasi database ke PostgreSQL untuk mendukung konkurensi yang lebih baik.",
    "Deploy ke platform cloud (Railway, Render, AWS) agar dapat diakses publik.",
    "Menambahkan fitur pencarian, kategori, dan timestamp untuk meningkatkan fungsionalitas.",
]
add_bullet_list(slide, Inches(7), Inches(2.15), Inches(5.5), Inches(4.0),
                suggestions, font_size=12)


# ==================================================
# SLIDE 16: DAFTAR PUSTAKA
# ==================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_light_bg(slide)
add_header_bar(slide, "Daftar Pustaka", "Referensi")

references = [
    '1.  Fielding, R. T. (2000). Architectural Styles and the Design of Network-based Software Architectures. Doctoral dissertation, UC Irvine.',
    '2.  Grinberg, M. (2018). Flask Web Development (2nd ed.). O\'Reilly Media.',
    '3.  Tanenbaum, A. S., & Van Steen, M. (2017). Distributed Systems (3rd ed.). Pearson.',
    '4.  Mozilla Developer Network. (2024). Fetch API. https://developer.mozilla.org/en-US/docs/Web/API/Fetch_API',
    '5.  Flask Documentation. (2024). Flask: A Python Microframework. https://flask.palletsprojects.com/',
    '6.  SQLite Documentation. (2024). SQLite: About. https://www.sqlite.org/about.html',
    '7.  Google. (2024). Google Keep API Reference. https://developers.google.com/keep/api',
    '8.  OWASP Foundation. (2023). OWASP Top Ten. https://owasp.org/www-project-top-ten/',
    '9.  Richardson, L., & Ruby, S. (2013). RESTful Web APIs. O\'Reilly Media.',
    '10. Pallets Projects. (2024). Jinja2 Documentation. https://jinja.palletsprojects.com/',
]

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True
for i, ref in enumerate(references):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = ref
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_BLACK
    p.font.name = "Calibri"
    p.space_after = Pt(8)

# Footer: Thank you
add_text_block(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
               "Terima Kasih — Apakah ada pertanyaan?",
               font_size=16, bold=True, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)


# ==================================================
# SAVE
# ==================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "Presentasi_Note_Taking_CSP.pptx")
prs.save(output_path)
print(f"\nPresentasi berhasil dibuat: {output_path}")
print(f"   Total slide: {len(prs.slides)}")
